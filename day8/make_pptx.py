# B-08-O1 — 모닝브루 사업 설명 PPTX 10장 덱 자동 생성
# 근거: PRD lite(B-02-O1)·DOCX 소개서(B-07-O2)·흐름도(B-07-O1)·시장조사(B-02-R1·R2)
# TAM 수치는 2026-07-15 웹 검색으로 확정 (검증 메모 B-08-O3 참조)
# 실행: .venv/bin/python day8/make_pptx.py
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BROWN = RGBColor(0x5C, 0x3D, 0x2E)
ACCENT = RGBColor(0xC4, 0x71, 0x3B)
CREAM = RGBColor(0xFA, 0xF8, 0xF5)
DIM = RGBColor(0x8A, 0x7A, 0x6E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x3A, 0x7A, 0x3A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height

def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = CREAM
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s

def txt(s, x, y, w, h, text, size=18, bold=False, color=BROWN, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
        r.font.name = '맑은 고딕'
    return box

def card(s, x, y, w, h, fill=WHITE, line_color=None):
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    c.fill.solid(); c.fill.fore_color.rgb = fill
    if line_color: c.line.color.rgb = line_color; c.line.width = Pt(1)
    else: c.line.fill.background()
    c.shadow.inherit = False
    return c

def header(s, no, title):
    txt(s, Inches(0.7), Inches(0.35), Inches(1.2), Inches(0.4), no, size=13, bold=True, color=ACCENT)
    txt(s, Inches(0.7), Inches(0.65), Inches(11), Inches(0.9), title, size=30, bold=True, color=BROWN)

# ═══ 01 표지 ═══
s = slide()
txt(s, 0, Inches(2.0), W, Inches(0.5), '사업 소개', size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
txt(s, 0, Inches(2.5), W, Inches(1.2), '모닝브루', size=60, bold=True, color=BROWN, align=PP_ALIGN.CENTER)
txt(s, 0, Inches(3.7), W, Inches(0.6), '"줄 서는 5분, 예약하면 0분."', size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
txt(s, 0, Inches(4.4), W, Inches(0.5), '망원동 핸드드립 카페 · 아침 픽업 예약 서비스', size=16, color=DIM, align=PP_ALIGN.CENTER)
txt(s, 0, Inches(5.6), W, Inches(0.8), '대표 정해린  ·  2026. 7. 15\n(바이브 코딩 기초반 실습용 가상 사업)', size=12, color=DIM, align=PP_ALIGN.CENTER)

# ═══ 02 문제 ═══
s = slide()
header(s, '02 문제', '출근길 아침, 커피 한 잔이 어렵다')
items = [
    ('⏰ 줄서기', '프랜차이즈 앞 대기 — 출근 시간에 쫓긴다'),
    ('❓ 예측 불가', '오늘 얼마나 기다릴지 알 수 없다'),
    ('☕ 품질 타협', '편의점 커피는 빠르지만 아쉽다'),
]
for i, (t1, t2) in enumerate(items):
    x = Inches(0.7 + i * 4.1)
    card(s, x, Inches(2.0), Inches(3.7), Inches(2.2))
    txt(s, x + Inches(0.3), Inches(2.3), Inches(3.1), Inches(0.6), t1, size=20, bold=True, color=ACCENT)
    txt(s, x + Inches(0.3), Inches(3.0), Inches(3.1), Inches(1.0), t2, size=15, color=BROWN)
txt(s, Inches(0.7), Inches(4.7), Inches(12), Inches(0.5),
    '타깃: 망원역 인근으로 출근하는 30대 직장인 — 아침 8~9시, 시간에 쫓기며 품질 좋은 커피를 원하는 사람', size=15, bold=True, color=BROWN)
txt(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.4),
    '※ 위 불편은 가설이며 검증을 진행 중입니다 (시장조사 B-02-R1 검증 항목)', size=11, color=DIM)

# ═══ 03 해결책 ═══
s = slide()
header(s, '03 해결책', '네 가지를 동시에 — 그래서 모닝브루')
values = [
    ('🌅 이른 아침', '매일 08시 오픈\n아침을 여는 카페'),
    ('☕ 핸드드립 품질', '한 잔씩 직접 내림\n당일 구운 스콘'),
    ('🪑 골목의 편안함', '망원동 골목 끝\n1인 좌석'),
    ('📝 사전 픽업 예약', '미리 주문하고\n픽업만 — 대기 0분'),
]
for i, (t1, t2) in enumerate(values):
    x = Inches(0.7 + i * 3.05)
    card(s, x, Inches(2.0), Inches(2.75), Inches(2.6))
    txt(s, x + Inches(0.25), Inches(2.3), Inches(2.25), Inches(0.9), t1, size=17, bold=True, color=ACCENT)
    txt(s, x + Inches(0.25), Inches(3.2), Inches(2.25), Inches(1.2), t2, size=13, color=BROWN)
txt(s, Inches(0.7), Inches(5.2), Inches(12), Inches(0.6),
    '"줄 서는 5분, 예약하면 0분." — 확인된 범위에서 네 가지를 함께 제공하는 대안은 드뭅니다 (상권 실사로 검증 예정)',
    size=14, bold=True, color=BROWN)

# ═══ 04 시장 규모 ★ (B-08-O2) ═══
s = slide()
header(s, '04 시장 규모 ★', 'TAM → SAM → SOM — 큰 시장에서 잡을 시장으로')
funnel = [
    (Inches(1.2), Inches(10.9), 'TAM · 국내 커피 시장', '약 10조 원 (2025)',
     '출처: 코리아비즈리뷰 2025 — 집계 기준별 10~21조 원 (Expert Market Research 외)', RGBColor(0xDD, 0xF4, 0xFF), RGBColor(0x05, 0x50, 0xAE)),
    (Inches(2.7), Inches(7.9), 'SAM · 망원역 상권 카페 시장', '약 100억 원 (직접 추정)',
     '산식: 상권 카페 약 50곳 × 평균 연매출 약 2억 원 — 상권 실사로 확정 필요', RGBColor(0xFD, 0xF8, 0xEE), RGBColor(0xA0, 0x60, 0x30)),
    (Inches(4.2), Inches(4.9), 'SOM · 아침 픽업 목표 (캐파 기반)', '월 약 515만 원',
     '산식: 일 최대 36잔(6슬롯×2시간×3잔) × 객단가 5,500원 × 월 26일 — 구현된 슬롯 정책 근거', RGBColor(0xEE, 0xF7, 0xEE), GREEN),
]
for x, w, t1, t2, src, fill, tc in funnel:
    y = Inches(1.75 + funnel.index((x, w, t1, t2, src, fill, tc)) * 1.55)
    c = card(s, x, y, w, Inches(1.35), fill=fill)
    txt(s, x + Inches(0.35), y + Inches(0.12), w - Inches(0.7), Inches(0.45), f'{t1}   —   {t2}', size=17, bold=True, color=tc)
    txt(s, x + Inches(0.35), y + Inches(0.62), w - Inches(0.7), Inches(0.6), src, size=11, color=DIM)
txt(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.5),
    '※ 모든 수치에 출처·산식을 표기 — TAM은 직접 검색으로 확정(검증 메모 B-08-O3), SAM·SOM은 추정임을 명시', size=11, color=DIM)

# ═══ 05 제품/서비스 ═══
s = slide()
header(s, '05 제품 / 서비스', '1분 예약, 0분 대기 — 실제로 동작하는 서비스')
steps = ['① 페이지 방문', '② 메뉴 확인', '③ 웹 폼 예약\n(날짜·시간·수량)', '④ 매장 픽업']
for i, st in enumerate(steps):
    x = Inches(0.7 + i * 1.85)
    card(s, x, Inches(1.9), Inches(1.6), Inches(1.1))
    txt(s, x + Inches(0.12), Inches(2.0), Inches(1.36), Inches(0.9), st, size=12, bold=True, color=BROWN, align=PP_ALIGN.CENTER)
ops = [
    '픽업: 08:00~10:00 · 10분 슬롯 (월 휴무)',
    '슬롯당 3잔 한정 — 마감 슬롯 선택 차단',
    '내 예약 확인·접수 취소 셀프 서비스',
    '접수→확정→픽업완료 상태 관리(관리자)',
    '개인정보 최소 수집 + 동의·보유기간 고지',
]
for i, op in enumerate(ops):
    txt(s, Inches(0.7), Inches(3.4 + i * 0.52), Inches(6.2), Inches(0.5), '·  ' + op, size=14, color=BROWN)
s.shapes.add_picture('day7/B-07-O1.png', Inches(7.4), Inches(1.9), height=Inches(4.9))
txt(s, Inches(7.4), Inches(6.85), Inches(5.2), Inches(0.4), '그림: 사용자 흐름도(오류 경로 포함)', size=10, color=DIM, align=PP_ALIGN.CENTER)
txt(s, Inches(0.7), Inches(6.4), Inches(6.2), Inches(0.5), '공개 서비스: vibe-coding-practice-xi.vercel.app/day1/', size=12, bold=True, color=ACCENT)

# ═══ 06 비즈니스 모델 ═══
s = slide()
header(s, '06 비즈니스 모델', '아침 픽업은 유입 장치 — 수익은 하루 전체에서')
cols = [
    ('수익원', ['음료·베이커리 판매 (객단가 4,500~9,000원)', '아침 픽업 → 단골 전환 → 종일 방문', '(확장 후보) 원두 판매·구독']),
    ('비용 구조', ['원재료(원두·유제품·베이커리)', '임대료 · 소모품', '1인 운영 — 인건비 최소화']),
    ('성장 논리', ['픽업 예약 = 재방문 습관화 장치', '슬롯 데이터로 수요 실측 → 캐파 조정', '검증 후 시간대·메뉴 확장']),
]
for i, (t1, items) in enumerate(cols):
    x = Inches(0.7 + i * 4.1)
    card(s, x, Inches(1.9), Inches(3.7), Inches(3.8))
    txt(s, x + Inches(0.3), Inches(2.15), Inches(3.1), Inches(0.5), t1, size=18, bold=True, color=ACCENT)
    for j, it in enumerate(items):
        txt(s, x + Inches(0.3), Inches(2.85 + j * 0.85), Inches(3.1), Inches(0.8), '· ' + it, size=13, color=BROWN)
txt(s, Inches(0.7), Inches(6.2), Inches(12), Inches(0.5),
    '※ 아침 캐파(월 ~515만 원)만으로는 손익 불충분 — 종일 매출과 결합하는 구조 (다면조사 투자자 검토 반영)', size=11, color=DIM)

# ═══ 07 창업자 ═══
s = slide()
header(s, '07 창업자', '매일 첫 잔을 직접 내리는 사장')
card(s, Inches(0.7), Inches(2.0), Inches(5.6), Inches(3.6))
txt(s, Inches(1.0), Inches(2.3), Inches(5.0), Inches(0.6), '정해린 · 모닝브루 대표', size=22, bold=True, color=BROWN)
for j, it in enumerate(['핸드드립을 내리는 1인 사장', '동네 손님의 아침 동선을 아는 운영자', '"품질을 지키는 만큼만 받는다" — 슬롯 3잔 원칙']):
    txt(s, Inches(1.0), Inches(3.1 + j * 0.65), Inches(5.0), Inches(0.6), '· ' + it, size=14, color=BROWN)
card(s, Inches(6.9), Inches(2.0), Inches(5.6), Inches(3.6))
txt(s, Inches(7.2), Inches(2.3), Inches(5.0), Inches(0.6), '운영 원칙', size=22, bold=True, color=ACCENT)
for j, it in enumerate(['매일 아침 직접 로스팅 원두 확인·추출', '당일 구운 스콘만 판매 (소진 시 사전 연락)', '개인정보 최소 수집 · 1개월 후 파기']):
    txt(s, Inches(7.2), Inches(3.1 + j * 0.65), Inches(5.0), Inches(0.6), '· ' + it, size=14, color=BROWN)
txt(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.4), '※ 실습용 가상 인물·프로필입니다', size=11, color=DIM)

# ═══ 08 트랙션 ═══
s = slide()
header(s, '08 트랙션', '이미 만든 것 · 앞으로 잴 것')
card(s, Inches(0.7), Inches(1.9), Inches(5.9), Inches(4.0))
txt(s, Inches(1.0), Inches(2.15), Inches(5.3), Inches(0.5), '✅ 지금까지 (사실)', size=18, bold=True, color=GREEN)
for j, it in enumerate(['공개 웹 서비스 배포 (Vercel, 자동 재배포)', '슬롯제 예약 + 관리자 시스템 가동', '자동 검증 19항목 체계 (기능·모바일)', 'SEO·GEO 기초 적용 (검색 친화 메타)']):
    txt(s, Inches(1.0), Inches(2.8 + j * 0.7), Inches(5.3), Inches(0.6), '· ' + it, size=13, color=BROWN)
card(s, Inches(6.9), Inches(1.9), Inches(5.7), Inches(4.0))
txt(s, Inches(7.2), Inches(2.15), Inches(5.1), Inches(0.5), '🎯 앞으로 잴 목표 지표', size=18, bold=True, color=ACCENT)
for j, it in enumerate(['아침 슬롯 가동률 50% (13슬롯×3잔 기준)', '노쇼율 10% 이하', '재방문(주 2회 이상) 고객 비중', '예약 → 픽업 완료 전환율']):
    txt(s, Inches(7.2), Inches(2.8 + j * 0.7), Inches(5.1), Inches(0.6), '· ' + it, size=13, color=BROWN)
txt(s, Inches(0.7), Inches(6.2), Inches(12), Inches(0.4), '※ 매출·이용 실적은 아직 없음 — 목표 지표로 제시 (베타 운영으로 실측 예정)', size=11, color=DIM)

# ═══ 09 요청 ═══
s = slide()
header(s, '09 요청', '함께해 주시면 빠르게 검증하겠습니다')
asks = [
    ('🙋 베타 고객 100명', '망원역 출근길 직장인 — 아침 픽업 예약 체험단'),
    ('📋 상권 실사 협조', 'SAM 추정치(상권 카페 수·매출) 확정을 위한 조사'),
    ('🤝 지역 파트너십', '인근 오피스·코워킹 — 아침 단체 픽업 제휴'),
]
for i, (t1, t2) in enumerate(asks):
    y = Inches(2.0 + i * 1.45)
    card(s, Inches(0.7), y, Inches(11.9), Inches(1.25))
    txt(s, Inches(1.1), y + Inches(0.15), Inches(3.6), Inches(0.9), t1, size=18, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(4.9), y + Inches(0.15), Inches(7.4), Inches(0.9), t2, size=14, color=BROWN, anchor=MSO_ANCHOR.MIDDLE)

# ═══ 10 마무리 ═══
s = slide()
txt(s, 0, Inches(2.2), W, Inches(0.8), '감사합니다', size=40, bold=True, color=BROWN, align=PP_ALIGN.CENTER)
txt(s, 0, Inches(3.2), W, Inches(0.5), '줄 서는 5분, 예약하면 0분 — 모닝브루', size=18, color=ACCENT, align=PP_ALIGN.CENTER)
txt(s, 0, Inches(4.2), W, Inches(1.6),
    '서울 망원동 골목 끝 (망원역 2번 출구 도보 6분)\n02-123-4567  ·  @morningbrew_official\nvibe-coding-practice-xi.vercel.app/day1/',
    size=14, color=DIM, align=PP_ALIGN.CENTER)

prs.save('day8/B-08-O1.pptx')
print(f'day8/B-08-O1.pptx 생성 완료 — 슬라이드 {len(prs.slides.__iter__.__self__._sldIdLst)}장')
